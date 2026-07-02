#!/usr/bin/env python3
"""Generate a diverse corpus of .pptx test files for Slideflow.

Self-contained: only depends on python-pptx (>= 1.0) and, optionally, Pillow
for image synthesis. Re-running is deterministic (fixed RNG seed) so the corpus
does not drift in size between runs.

Usage:
    python3 scripts/generate_examples.py            # (re)generate every deck
    python3 scripts/generate_examples.py tables     # generate only tables.pptx
    python3 scripts/generate_examples.py --list     # list deck names

Output goes to examples/pptx/ next to this repo. Downloaded real-world decks
live in examples/pptx/real/ and are NOT touched by this script.
"""
from __future__ import annotations

import io
import math
import random
import sys
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

try:
    from PIL import Image, ImageDraw

    HAVE_PIL = True
except ImportError:  # pragma: no cover - Pillow is optional
    HAVE_PIL = False

REPO_ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = REPO_ROOT / "examples" / "pptx"

# 16:9 widescreen (default for most decks) and classic 4:3.
WIDE = (Emu(12192000), Emu(6858000))
CLASSIC_4X3 = (Emu(9144000), Emu(6858000))

RNG = random.Random(1234)

# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------


def new_prs(size=WIDE) -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = size
    return prs


def set_props(prs: Presentation, title: str, author: str = "Slideflow Corpus Generator") -> None:
    """docProps title deliberately DIFFERENT from the file name so the app's
    'show real file name, not docProps title' behaviour is exercised."""
    cp = prs.core_properties
    cp.title = title
    cp.author = author
    cp.category = "Slideflow test corpus"


def add_title_only(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = title
    return slide


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    return tb.text_frame


def _png_bytes(img) -> bytes:
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _jpeg_bytes(img, quality=82) -> bytes:
    buf = io.BytesIO()
    img.convert("RGB").save(buf, format="JPEG", quality=quality)
    return buf.getvalue()


# --- Image synthesis (Pillow if present, else hand-rolled raw PNG) ----------

# 1x1 transparent PNG, valid file, used when we need a truly tiny image and as
# a fallback if Pillow is missing.
_PNG_1X1 = bytes(
    [
        0x89, 0x50, 0x4E, 0x47, 0x0D, 0x0A, 0x1A, 0x0A, 0x00, 0x00, 0x00, 0x0D,
        0x49, 0x48, 0x44, 0x52, 0x00, 0x00, 0x00, 0x01, 0x00, 0x00, 0x00, 0x01,
        0x08, 0x06, 0x00, 0x00, 0x00, 0x1F, 0x15, 0xC4, 0x89, 0x00, 0x00, 0x00,
        0x0D, 0x49, 0x44, 0x41, 0x54, 0x78, 0x9C, 0x62, 0x00, 0x01, 0x00, 0x00,
        0x05, 0x00, 0x01, 0x0D, 0x0A, 0x2D, 0xB4, 0x00, 0x00, 0x00, 0x00, 0x49,
        0x45, 0x4E, 0x44, 0xAE, 0x42, 0x60, 0x82,
    ]
)


def img_transparent_png(size=256) -> bytes:
    """PNG with a real alpha channel: a translucent disc on transparency."""
    if not HAVE_PIL:
        return _PNG_1X1
    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    d.ellipse([16, 16, size - 16, size - 16], fill=(70, 130, 220, 180))
    d.ellipse([size // 3, size // 3, size - 40, size - 40], fill=(240, 180, 40, 210))
    d.line([0, 0, size, size], fill=(200, 40, 60, 255), width=8)
    return _png_bytes(img)


def img_gradient_jpeg(w=640, h=400) -> bytes:
    """Photo-ish smooth gradient with noise, saved as JPEG."""
    if not HAVE_PIL:
        return _PNG_1X1
    img = Image.new("RGB", (w, h))
    px = img.load()
    for y in range(h):
        for x in range(w):
            r = int(40 + 200 * x / w)
            g = int(60 + 160 * y / h)
            b = int(120 + 100 * math.sin(x / 40.0) * math.cos(y / 30.0))
            b = max(0, min(255, b))
            px[x, y] = (r, g, b)
    # sprinkle a little noise so it isn't a perfectly compressible gradient
    d = ImageDraw.Draw(img)
    for _ in range(300):
        x, y = RNG.randint(0, w - 1), RNG.randint(0, h - 1)
        d.point((x, y), fill=(RNG.randint(0, 255), RNG.randint(0, 255), RNG.randint(0, 255)))
    return _jpeg_bytes(img)


def img_logo_png(size=200) -> bytes:
    """A small deterministic 'logo' reused across slides/decks to test dedup."""
    if not HAVE_PIL:
        return _PNG_1X1
    img = Image.new("RGBA", (size, size), (255, 255, 255, 0))
    d = ImageDraw.Draw(img)
    d.rounded_rectangle([10, 10, size - 10, size - 10], radius=24, fill=(30, 40, 90, 255))
    d.polygon([(size // 2, 30), (size - 40, size - 40), (40, size - 40)], fill=(250, 210, 60, 255))
    d.ellipse([size // 2 - 22, size // 2 - 22, size // 2 + 22, size // 2 + 22], fill=(220, 60, 80, 255))
    return _png_bytes(img)


def pic_from_bytes(slide, data: bytes, left, top, width=None, height=None):
    return slide.shapes.add_picture(io.BytesIO(data), left, top, width=width, height=height)


# ---------------------------------------------------------------------------
# Deck builders
# ---------------------------------------------------------------------------


def build_text_hierarchy(path: Path) -> None:
    prs = new_prs()
    set_props(prs, "Generated Example — Text Hierarchy & Typography")

    # 1. Title slide
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Texthierarchie & Typografie"
    s.placeholders[1].text = "Ein Beispiel-Deck mit vielen Ebenen, Fonts und Farben"

    # 2. Section header
    s = prs.slides.add_slide(prs.slide_layouts[2])
    s.shapes.title.text = "Abschnitt 1 — Aufzählungsebenen"
    s.placeholders[1].text = "Fünf Ebenen tiefer Verschachtelung"

    # 3. Five-level bullets
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Fünf Aufzählungsebenen"
    tf = s.placeholders[1].text_frame
    tf.text = "Ebene 1 — oberste Ebene"
    for lvl in range(1, 5):
        p = tf.add_paragraph()
        p.text = f"Ebene {lvl + 1} — verschachtelt auf Tiefe {lvl}"
        p.level = lvl

    # 4. Explicit fonts / sizes / colors / bold / italic
    s = add_title_only(prs, "Explizite Font-Formatierung")
    tf = textbox(s, Inches(0.7), Inches(1.6), Inches(11), Inches(4.5))
    specs = [
        ("Große fette Überschrift", 40, True, False, RGBColor(0x1F, 0x3A, 0x93), "Arial Black"),
        ("Kursiver Untertitel in Grün", 24, False, True, RGBColor(0x1E, 0x7A, 0x33), "Georgia"),
        ("Kleiner grauer Fließtext", 14, False, False, RGBColor(0x66, 0x66, 0x66), "Calibri"),
        ("Roter Warnhinweis, fett + kursiv", 20, True, True, RGBColor(0xC0, 0x28, 0x28), "Verdana"),
    ]
    for i, (text, size, bold, italic, color, font) in enumerate(specs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font

    # 5. Mixed runs within one paragraph + explicit line break
    s = add_title_only(prs, "Gemischte Runs & Zeilenumbruch")
    tf = textbox(s, Inches(0.7), Inches(1.8), Inches(11), Inches(3))
    p = tf.paragraphs[0]
    for text, bold, color in [
        ("Dieser Satz ", False, RGBColor(0, 0, 0)),
        ("enthält ", True, RGBColor(0xB0, 0x30, 0x10)),
        ("mehrere ", False, RGBColor(0x10, 0x60, 0xB0)),
        ("formatierte Runs.", True, RGBColor(0x20, 0x80, 0x20)),
    ]:
        r = p.add_run()
        r.text = text
        r.font.bold = bold
        r.font.color.rgb = color
    # explicit line break (<a:br>) between two runs in the same paragraph
    ra = p.add_run()
    ra.text = "Zeile A"
    p._p.append(p._p.makeelement(qn("a:br"), {}))
    rb = p.add_run()
    rb.text = "Zeile B nach Umbruch"

    # 6. Long paragraph (wrapping / autofit stress)
    s = add_title_only(prs, "Langer Fließtext")
    tf = textbox(s, Inches(0.7), Inches(1.6), Inches(11.5), Inches(5))
    tf.word_wrap = True
    long_para = (
        "Lorem ipsum dolor sit amet, consetetur sadipscing elitr, sed diam nonumy "
        "eirmod tempor invidunt ut labore et dolore magna aliquyam erat, sed diam "
        "voluptua. At vero eos et accusam et justo duo dolores et ea rebum. Stet clita "
        "kasd gubergren, no sea takimata sanctus est Lorem ipsum dolor sit amet. "
    ) * 3
    tf.text = long_para

    # 7. Alignment showcase
    s = add_title_only(prs, "Absatz-Ausrichtung")
    tf = textbox(s, Inches(0.7), Inches(1.6), Inches(11.5), Inches(4.5))
    for align, label in [
        (PP_ALIGN.LEFT, "Linksbündig"),
        (PP_ALIGN.CENTER, "Zentriert"),
        (PP_ALIGN.RIGHT, "Rechtsbündig"),
        (PP_ALIGN.JUSTIFY, "Blocksatz " + "Wort " * 20),
    ]:
        p = tf.paragraphs[0] if tf.paragraphs[0].text == "" and align == PP_ALIGN.LEFT else tf.add_paragraph()
        p.text = label
        p.alignment = align

    # 8. Numbered-ish / dense outline
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Dichte Gliederung"
    tf = s.placeholders[1].text_frame
    tf.text = "Punkt eins mit etwas mehr Text zur Zeile"
    for i in range(2, 12):
        p = tf.add_paragraph()
        p.text = f"Punkt {i} — Detailzeile mit Zusatzinformationen"
        p.level = i % 3

    prs.save(str(path))


def _header_cell(cell, text, fill=RGBColor(0x1F, 0x3A, 0x93)):
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    run = cell.text_frame.paragraphs[0].runs[0]
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.bold = True


def build_tables(path: Path) -> None:
    prs = new_prs()
    set_props(prs, "Generated Example — Tables & Grids")

    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Tabellen & Raster"
    s.placeholders[1].text = "Kopfzeilen, verbundene Zellen, Zellfarben, Spaltenbreiten"

    # Header row + zebra fills
    s = add_title_only(prs, "Tabelle mit Kopfzeile & Zebra-Füllung")
    tbl = s.shapes.add_table(4, 4, Inches(0.6), Inches(1.6), Inches(11), Inches(3.5)).table
    for c, h in enumerate(["Quartal", "Umsatz", "Kosten", "Gewinn"]):
        _header_cell(tbl.cell(0, c), h)
    body = [["Q1", "120 T€", "80 T€", "40 T€"],
            ["Q2", "150 T€", "90 T€", "60 T€"],
            ["Q3", "170 T€", "95 T€", "75 T€"]]
    for r, row in enumerate(body, start=1):
        for c, v in enumerate(row):
            tbl.cell(r, c).text = v
            if r % 2 == 0:
                tbl.cell(r, c).fill.solid()
                tbl.cell(r, c).fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xFB)

    # Merged cells (horizontal + vertical) + custom column widths
    s = add_title_only(prs, "Verbundene Zellen & Spaltenbreiten")
    tbl = s.shapes.add_table(5, 4, Inches(0.6), Inches(1.6), Inches(11), Inches(3.8)).table
    for i, w in enumerate([Inches(4), Inches(2.6), Inches(2.4), Inches(2)]):
        tbl.columns[i].width = w
    tbl.cell(0, 0).merge(tbl.cell(0, 1))
    _header_cell(tbl.cell(0, 0), "Verbundener Kopf (2 Spalten)")
    _header_cell(tbl.cell(0, 2), "A")
    _header_cell(tbl.cell(0, 3), "B")
    tbl.cell(1, 0).merge(tbl.cell(2, 0))
    tbl.cell(1, 0).text = "Vertikal verbunden"
    tbl.cell(1, 0).vertical_anchor = MSO_ANCHOR.MIDDLE
    for r in range(1, 5):
        for c in range(4):
            if not tbl.cell(r, c).text:
                tbl.cell(r, c).text = f"R{r}·C{c}"

    # Large 10x8 table
    s = add_title_only(prs, "Große Tabelle 10×8")
    tbl = s.shapes.add_table(10, 8, Inches(0.6), Inches(1.4), Inches(12), Inches(5.3)).table
    for c in range(8):
        _header_cell(tbl.cell(0, c), f"Sp {c + 1}", fill=RGBColor(0x33, 0x33, 0x33))
    for r in range(1, 10):
        for c in range(8):
            tbl.cell(r, c).text = str(RNG.randint(0, 999))

    prs.save(str(path))


def build_charts(path: Path) -> None:
    prs = new_prs()
    set_props(prs, "Generated Example — Charts & Embedded Workbooks")

    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Diagramme"
    s.placeholders[1].text = "Balken, gestapelt, Linie, Kreis, Ring, Streuung — je mit eingebettetem xlsx"

    cats = ["Jan", "Feb", "Mär", "Apr"]

    def cat_data(series):
        cd = CategoryChartData()
        cd.categories = cats
        for name, vals in series:
            cd.add_series(name, vals)
        return cd

    def chart_slide(title, ctype, cd):
        sl = add_title_only(prs, title)
        gf = sl.shapes.add_chart(ctype, Inches(1), Inches(1.6), Inches(11), Inches(5), cd)
        return gf.chart

    c = chart_slide("Balkendiagramm (gruppiert)", XL_CHART_TYPE.COLUMN_CLUSTERED,
                    cat_data([("Umsatz", (19, 21, 17, 24)), ("Kosten", (12, 14, 11, 15))]))
    c.has_legend = True
    c.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart_slide("Gestapelte Balken", XL_CHART_TYPE.COLUMN_STACKED,
                cat_data([("A", (5, 6, 7, 8)), ("B", (3, 2, 4, 1)), ("C", (2, 3, 1, 2))]))
    chart_slide("Liniendiagramm", XL_CHART_TYPE.LINE,
                cat_data([("Serie 1", (10, 14, 9, 17)), ("Serie 2", (7, 8, 12, 6))]))
    c = chart_slide("Kreisdiagramm", XL_CHART_TYPE.PIE, cat_data([("Anteile", (35, 25, 20, 20))]))
    c.plots[0].has_data_labels = True
    chart_slide("Ringdiagramm", XL_CHART_TYPE.DOUGHNUT, cat_data([("Ring", (40, 30, 15, 15))]))

    s = add_title_only(prs, "Streudiagramm (XY)")
    xy = XyChartData()
    ser = xy.add_series("Punkte")
    for _ in range(12):
        ser.add_data_point(RNG.uniform(0, 10), RNG.uniform(0, 10))
    s.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, Inches(1), Inches(1.6), Inches(11), Inches(5), xy)

    prs.save(str(path))


def build_images(path: Path) -> None:
    prs = new_prs()
    set_props(prs, "Generated Example — Images & Media Reuse")
    logo, trans, jpeg = img_logo_png(), img_transparent_png(), img_gradient_jpeg()

    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Bilder & Medien"
    s.placeholders[1].text = "Transparenz, JPEG-Foto, Wiederverwendung (Dedup), Rotation, 1×1"

    s = add_title_only(prs, "PNG mit Transparenz")
    pic_from_bytes(s, trans, Inches(4.5), Inches(1.7), height=Inches(4))

    s = add_title_only(prs, "JPEG-Foto (Verlauf)")
    pic_from_bytes(s, jpeg, Inches(2.5), Inches(1.7), width=Inches(8))

    # Same logo reused across 3 slides (twice each) → dedup test
    for i in range(3):
        s = add_title_only(prs, f"Wiederverwendetes Logo #{i + 1}")
        pic_from_bytes(s, logo, Inches(4.5), Inches(2), height=Inches(2.5))
        pic_from_bytes(s, logo, Inches(8), Inches(3.6), height=Inches(1.5))

    s = add_title_only(prs, "Rotiertes Bild")
    p = pic_from_bytes(s, logo, Inches(4.8), Inches(2), height=Inches(3))
    p.rotation = 30

    s = add_title_only(prs, "Winziges 1×1-Bild")
    pic_from_bytes(s, _PNG_1X1, Inches(6.4), Inches(3.4), width=Inches(0.25), height=Inches(0.25))

    prs.save(str(path))


def _add_outer_shadow(shape) -> None:
    spPr = shape._element.spPr
    eff = spPr.makeelement(qn("a:effectLst"), {})
    shd = spPr.makeelement(qn("a:outerShdw"),
                           {"blurRad": "40000", "dist": "23000", "dir": "5400000", "rotWithShape": "0"})
    clr = spPr.makeelement(qn("a:srgbClr"), {"val": "000000"})
    clr.append(spPr.makeelement(qn("a:alpha"), {"val": "38000"}))
    shd.append(clr)
    eff.append(shd)
    spPr.append(eff)


def build_shapes(path: Path) -> None:
    prs = new_prs()
    set_props(prs, "Generated Example — Shapes, Connectors & Freeform")

    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Formen & Verbinder"
    s.placeholders[1].text = "Autoshapes, Pfeile, Sterne, Chevrons, Gruppen, Freiform"

    # Autoshape gallery
    s = add_title_only(prs, "Autoshapes")
    gallery = [
        (MSO_SHAPE.ROUNDED_RECTANGLE, RGBColor(0x2E, 0x6D, 0xB5)),
        (MSO_SHAPE.RIGHT_ARROW, RGBColor(0x1E, 0x9E, 0x5A)),
        (MSO_SHAPE.STAR_5_POINT, RGBColor(0xE0, 0xA1, 0x1A)),
        (MSO_SHAPE.CHEVRON, RGBColor(0xB0, 0x3A, 0x2E)),
        (MSO_SHAPE.PENTAGON, RGBColor(0x6A, 0x3D, 0x9A)),
        (MSO_SHAPE.OVAL, RGBColor(0x2A, 0x9D, 0x8F)),
    ]
    x = Inches(0.5)
    for shp, color in gallery:
        sp = s.shapes.add_shape(shp, x, Inches(2.2), Inches(1.8), Inches(1.5))
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
        sp.line.color.rgb = RGBColor(0x22, 0x22, 0x22)
        sp.line.width = Pt(2)
        x += Inches(2.05)

    # Gradient / solid / dashed outline / explicit shadow
    s = add_title_only(prs, "Füllungen, Umrisse, Schatten")
    grad = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2), Inches(3.4), Inches(2.6))
    grad.fill.gradient()
    dash = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.7), Inches(2), Inches(3.4), Inches(2.6))
    dash.fill.solid()
    dash.fill.fore_color.rgb = RGBColor(0xD9, 0x53, 0x4F)
    dash.line.color.rgb = RGBColor(0x22, 0x22, 0x22)
    dash.line.width = Pt(3)
    dash.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    shadowed = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.6), Inches(2), Inches(3.4), Inches(2.6))
    shadowed.fill.solid()
    shadowed.fill.fore_color.rgb = RGBColor(0x3B, 0x7D, 0xC4)
    _add_outer_shadow(shadowed)

    # Connectors
    s = add_title_only(prs, "Verbinder")
    s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(2), Inches(1.8), Inches(1.3))
    s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(4), Inches(1.8), Inches(1.3))
    elbow = s.shapes.add_connector(MSO_CONNECTOR.ELBOW, Inches(2.8), Inches(2.65), Inches(9.2), Inches(4.65))
    elbow.line.width = Pt(2.25)
    s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(5.6), Inches(11), Inches(5.6))

    # Group shape (two rects + a connector inside a group)
    s = add_title_only(prs, "Gruppierte Formen")
    grp = s.shapes.add_group_shape()
    a = grp.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2.2), Inches(2.4), Inches(1.3))
    a.fill.solid()
    a.fill.fore_color.rgb = RGBColor(0x2E, 0x6D, 0xB5)
    b = grp.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.4), Inches(3.4), Inches(2.4), Inches(1.3))
    b.fill.solid()
    b.fill.fore_color.rgb = RGBColor(0x1E, 0x9E, 0x5A)
    grp.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.4), Inches(2.85), Inches(5.4), Inches(4.05))

    # Freeform polygon (star-ish path)
    s = add_title_only(prs, "Freiform-Polygon")
    fb = s.shapes.build_freeform(60, 10, scale=Pt(2.4))
    pts = [(75, 45), (115, 45), (82, 68), (95, 105), (60, 82), (25, 105), (38, 68), (5, 45), (45, 45)]
    fb.add_line_segments(pts, close=True)
    free = fb.convert_to_shape(origin_x=Inches(4.5), origin_y=Inches(1.8))
    free.fill.solid()
    free.fill.fore_color.rgb = RGBColor(0xE0, 0xA1, 0x1A)

    prs.save(str(path))


def build_notes_and_links(path: Path) -> None:
    prs = new_prs()
    set_props(prs, "Generated Example — Notes, Hyperlinks & Jumps")

    first = prs.slides.add_slide(prs.slide_layouts[0])
    first.shapes.title.text = "Notizen & Hyperlinks"
    first.placeholders[1].text = "Sprechernotizen mit Umlauten, externe Links, interner Sprung"

    # Multi-paragraph speaker notes with umlauts
    s = add_title_only(prs, "Folie mit Sprechernotizen")
    notes = s.notes_slide.notes_text_frame
    notes.text = "Erster Absatz der Notizen — mit Umlauten: ä ö ü Ä Ö Ü ß."
    for para in ["Zweiter Absatz: Grüße an das Prüf-Team!",
                 "Dritter Absatz mit Aufzählungszeichen: • Punkt A • Punkt B • Straße"]:
        notes.add_paragraph().text = para

    # External hyperlinks on text runs
    s = add_title_only(prs, "Externe Hyperlinks")
    tf = textbox(s, Inches(1), Inches(2), Inches(10.5), Inches(3))
    for i, (label, url) in enumerate([
        ("python-pptx Dokumentation", "https://python-pptx.readthedocs.io/"),
        ("OOXML-Spezifikation", "https://learn.microsoft.com/openspecs/office_standards/"),
        ("Beispiel-Domain mit Umlaut-Query", "https://example.com/pfad?stadt=zürich"),
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = f"→ {label}"
        r.hyperlink.address = url

    # Internal slide-to-slide jump via click_action.target_slide
    s = add_title_only(prs, "Interner Sprung zur ersten Folie")
    jump = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(3), Inches(6), Inches(1.3))
    jump.fill.solid()
    jump.fill.fore_color.rgb = RGBColor(0x2E, 0x6D, 0xB5)
    jump.text_frame.text = "Springe zur ersten Folie"
    jump.click_action.target_slide = first

    prs.save(str(path))


def build_classic_4x3(path: Path) -> None:
    prs = new_prs(CLASSIC_4X3)
    set_props(prs, "Generated Example — Classic 4:3 Layout")

    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Klassisches 4:3-Format"
    s.placeholders[1].text = "9144000 × 6858000 EMU — für gemischte Kompositionsgrößen"

    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Inhalt im 4:3-Format"
    tf = s.placeholders[1].text_frame
    tf.text = "Erster Punkt"
    for i in range(2, 5):
        tf.add_paragraph().text = f"Punkt {i} im 4:3-Raster"

    s = add_title_only(prs, "Bild im 4:3-Format")
    pic_from_bytes(s, img_logo_png(), Inches(3), Inches(2), height=Inches(3))

    prs.save(str(path))


def build_unicode_i18n(path: Path) -> None:
    prs = new_prs()
    set_props(prs, "Generated Example — Unicode & Internationalisierung")

    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Unicode & i18n"
    s.placeholders[1].text = "Umlaute, CJK, Kyrillisch, Griechisch, Emoji, RTL, Diakritika"

    samples = [
        ("Deutsch", "Grüße, Straße, Fußgängerübergänge — ä ö ü Ä Ö Ü ß"),
        ("中文/日本語/한국어", "你好世界，这是测试。こんにちは。안녕하세요 세계"),
        ("Кириллица", "Привет, мир! Это тестовый слайд на русском языке."),
        ("Ελληνικά", "Γειά σου Κόσμε! Δοκιμαστική διαφάνεια στα ελληνικά."),
        ("Emoji", "🎉 🚀 ✅ 📊 🖼️ 🇩🇪 🇯🇵 — Farben & Symbole"),
        ("العربية (RTL)", "مرحبا بالعالم! هذه شريحة اختبار باللغة العربية."),
        ("עברית (RTL)", "שלום עולם! זוהי שקופית בדיקה בעברית."),
        ("Diakritika", "café naïve façade Zürich piñata čačić łódź ñ ç ș ț đ"),
    ]
    s = add_title_only(prs, "Mehrsprachige Textzeilen")
    tf = textbox(s, Inches(0.6), Inches(1.5), Inches(12), Inches(5.2))
    tf.word_wrap = True
    for i, (lang, text) in enumerate(samples):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        head = p.add_run()
        head.text = f"{lang}: "
        head.font.bold = True
        p.add_run().text = text

    # Title packed with searchable diacritics + scripts (search/index test)
    add_title_only(prs, "Suchtest: Zürich Café — 東京 — Москва — Αθήνα — 🎯")

    prs.save(str(path))


def build_big_deck(path: Path) -> None:
    prs = new_prs()
    set_props(prs, "Generated Example — Large 40-Slide Deck")
    logo = img_logo_png()
    layout_cycle = [0, 1, 2, 5, 1, 3, 1, 5]
    topics = ["Einführung", "Marktanalyse", "Strategie", "Produkt", "Technik",
              "Team", "Finanzen", "Roadmap", "Risiken", "Zusammenfassung"]
    for i in range(40):
        s = prs.slides.add_slide(prs.slide_layouts[layout_cycle[i % len(layout_cycle)]])
        topic = topics[i % len(topics)]
        if s.shapes.title is not None:
            s.shapes.title.text = f"Folie {i + 1}: {topic}"
        for ph in s.placeholders:
            if ph.placeholder_format.idx == 1 and ph.has_text_frame:
                tf = ph.text_frame
                tf.text = f"Kernaussage {i + 1}"
                for j in range(1, 4):
                    p = tf.add_paragraph()
                    p.text = f"Detail {j} zu {topic}"
                    p.level = j % 2
                break
        if i % 5 == 0:
            pic_from_bytes(s, logo, Inches(11.2), Inches(0.25), height=Inches(0.8))
    prs.save(str(path))


def build_edge_cases(path: Path) -> None:
    prs = new_prs()
    set_props(prs, "Generated Example — Edge Cases")

    # Empty slide (blank layout, no shapes at all)
    prs.slides.add_slide(prs.slide_layouts[6])

    # Slide with only an image, no text
    s = prs.slides.add_slide(prs.slide_layouts[6])
    pic_from_bytes(s, img_gradient_jpeg(), Inches(2), Inches(1.4), width=Inches(9))

    # Very long title (300+ chars)
    add_title_only(prs, "Sehr langer Titel: " + "Wort " * 70)

    # Deeply nested groups (4 levels) around one leaf shape
    s = add_title_only(prs, "Tief verschachtelte Gruppe")
    g = s.shapes
    for _ in range(4):
        g = g.add_group_shape().shapes
    leaf = g.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(5), Inches(3), Inches(2.2), Inches(2.2))
    leaf.fill.solid()
    leaf.fill.fore_color.rgb = RGBColor(0xE0, 0xA1, 0x1A)

    # Text box with shrink-to-fit autofit
    s = add_title_only(prs, "Textfeld mit Autofit")
    tf = textbox(s, Inches(1), Inches(2), Inches(6), Inches(2))
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.text = "Dieser Text soll automatisch an die Form angepasst werden. " * 4

    prs.save(str(path))


# ---------------------------------------------------------------------------
# Registry + CLI
# ---------------------------------------------------------------------------

BUILDERS = {
    "text-hierarchy": build_text_hierarchy,
    "tables": build_tables,
    "charts": build_charts,
    "images": build_images,
    "shapes": build_shapes,
    "notes-and-links": build_notes_and_links,
    "classic-4x3": build_classic_4x3,
    "unicode-i18n": build_unicode_i18n,
    "big-deck": build_big_deck,
    "edge-cases": build_edge_cases,
}


def main(argv: list[str]) -> int:
    if "--list" in argv:
        for name in BUILDERS:
            print(name)
        return 0
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    wanted = [a for a in argv if not a.startswith("-")]
    names = wanted if wanted else list(BUILDERS)
    if not HAVE_PIL:
        print("WARNING: Pillow not installed — image decks will use 1x1 placeholders.", file=sys.stderr)
    for name in names:
        if name not in BUILDERS:
            print(f"unknown deck: {name} (known: {', '.join(BUILDERS)})", file=sys.stderr)
            return 2
        path = OUT_DIR / f"{name}.pptx"
        BUILDERS[name](path)
        size_kb = path.stat().st_size / 1024
        print(f"wrote {path.relative_to(REPO_ROOT)}  ({size_kb:.0f} KB)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
